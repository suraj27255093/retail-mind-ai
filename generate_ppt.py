import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx not yet ready")
    sys.exit(1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
LIGHT_BLUE = RGBColor(147, 197, 253)
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(30, 41, 59)
MUTED_TEXT = RGBColor(100, 116, 139)
CARD_BG = RGBColor(248, 250, 252)
EMERALD = RGBColor(16, 185, 129)

def add_header(slide, category, title):
    # Category badge
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.4))
    tf = cat_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = category.upper()
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.8))
    tf2 = title_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = NAVY

# ── SLIDE 1: TITLE SLIDE ─────────────────────────────
slide1 = prs.slides.add_slide(blank_layout)
bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
bg1.fill.solid()
bg1.fill.fore_color.rgb = NAVY
bg1.line.color.rgb = NAVY

t_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
tf1 = t_box.text_frame
tf1.word_wrap = True

p1 = tf1.paragraphs[0]
p1.text = "🛒 RetailMind AI"
p1.font.size = Pt(48)
p1.font.bold = True
p1.font.color.rgb = WHITE
p1.alignment = PP_ALIGN.LEFT

p2 = tf1.add_paragraph()
p2.text = "Smart Enterprise Retail & Mandi Price Intelligence System"
p2.font.size = Pt(22)
p2.font.color.rgb = LIGHT_BLUE
p2.space_before = Pt(14)

p3 = tf1.add_paragraph()
p3.text = "AI-Powered Solution for Indian Grocery Stores, Supermarkets & Wholesale Distributors"
p3.font.size = Pt(14)
p3.font.color.rgb = RGBColor(203, 213, 225)
p3.space_before = Pt(20)

p4 = tf1.add_paragraph()
p4.text = "Architected & Developed by: Suraj V. Shewale  |  Enterprise v3.0 Active"
p4.font.size = Pt(13)
p4.font.bold = True
p4.font.color.rgb = EMERALD
p4.space_before = Pt(30)

# ── SLIDE 2: EXECUTIVE SUMMARY ───────────────────────
slide2 = prs.slides.add_slide(blank_layout)
add_header(slide2, "Executive Overview", "RetailMind AI — Transforming Indian Commercial Retail")

card_w, card_h = Inches(3.6), Inches(4.8)
card_data = [
    ("🌾 Mandi Price Intelligence", "Priority 1 Government Data", "Integrates live APMC Mandi feeds (Agmarknet & MSAMB) for Nashik, Pune & Malegaon wholesale markets with multi-price attribute tracking.", BLUE),
    ("📦 ML Inventory Radar", "Demand & Expiry Forecasting", "Automated Reorder Point calculation (ROP = d × L + SS) preventing stockouts and flagging expiry risks 45 days in advance.", EMERALD),
    ("🧾 Ultra-Fast POS Billing", "Instant Customer Engagement", "Barcode cart calculation, GST compliance, custom discounts, and 1-click WhatsApp customer digital receipts in < 1 sec.", NAVY)
]

for i, (ctitle, csub, cdesc, accent) in enumerate(card_data):
    left = Inches(0.8 + i * 3.9)
    card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = RGBColor(226, 232, 240)
    
    # Accent top border
    border = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.8), card_w, Inches(0.12))
    border.fill.solid()
    border.fill.fore_color.rgb = accent
    border.line.fill.background()
    
    tb = slide2.shapes.add_textbox(left + Inches(0.2), Inches(2.1), card_w - Inches(0.4), card_h - Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = ctitle
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    p_sub = tf.add_paragraph()
    p_sub.text = csub
    p_sub.font.size = Pt(12)
    p_sub.font.bold = True
    p_sub.font.color.rgb = accent
    p_sub.space_before = Pt(6)
    
    p_desc = tf.add_paragraph()
    p_desc.text = cdesc
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = DARK_TEXT
    p_desc.space_before = Pt(14)

# ── SLIDE 3: PROBLEM STATEMENT ───────────────────────
slide3 = prs.slides.add_slide(blank_layout)
add_header(slide3, "Market Dynamics", "The Core Challenges in Traditional Grocery & Wholesale Retail")

probs = [
    ("❌ Unverified Sourcing Rates", "Retailers buy from wholesale markets without knowing real-time APMC Mandi rates, resulting in 8-15% margin leakage."),
    ("❌ Manual Stockout & Expiry Management", "Overstocking perishable items leads to 12% wastage, while unexpected stockouts cause customer drop-offs."),
    ("❌ Cluttered POS & Slow Receipts", "Legacy billing software is slow, lacks WhatsApp integration, and fails to capture customer loyalty metrics.")
]

for i, (title, desc) in enumerate(probs):
    top = Inches(1.8 + i * 1.7)
    box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.733), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(254, 242, 242)
    box.line.color.rgb = RGBColor(254, 202, 202)
    
    tb = slide3.shapes.add_textbox(Inches(1.1), top + Inches(0.15), Inches(11.1), Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = RGBColor(185, 28, 28)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = DARK_TEXT
    p2.space_before = Pt(4)

# ── SLIDE 4: MANDI INTELLIGENCE ──────────────────────
slide4 = prs.slides.add_slide(blank_layout)
add_header(slide4, "Government Feeds", "Priority 1 Government APMC Wholesale Mandi Intelligence")

tb4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0))
tf4 = tb4.text_frame
tf4.word_wrap = True

pts4 = [
    ("🏛️ Priority 1 Data Sources:", "Connects directly to official Agmarknet, eNAM, and MSAMB Mandi portals (Nashik, Pune, Malegaon). Never uses retail e-commerce site prices."),
    ("📊 Multi-Price Attribute Schema:", "Tracks Purchase Rate (Wholesale), Wholesale Average, Retail MRP, and 7-Day Market Average with Government Source Attribution & Confidence Scores."),
    ("🔄 Dynamic Real-time Rate Sync:", "Instant live price refresh engine capturing daily arrival fluctuations (-2.5% to +3.5%) to optimize procurement timing."),
    ("⚠️ Offline Price Fallback:", "Explicitly displays 'Live market price unavailable. Showing last verified market price' when live feeds undergo government portal maintenance.")
]

for h, b in pts4:
    p = tf4.add_paragraph()
    p.text = h
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.space_before = Pt(12)
    
    p_b = tf4.add_paragraph()
    p_b.text = b
    p_b.font.size = Pt(14)
    p_b.font.color.rgb = DARK_TEXT
    p_b.space_before = Pt(4)

# ── SLIDE 5: ML FORECASTING ──────────────────────────
slide5 = prs.slides.add_slide(blank_layout)
add_header(slide5, "Machine Learning Radar", "ML-Driven Inventory Demand Forecasting & Expiry Watchlist")

tb5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0))
tf5 = tb5.text_frame
tf5.word_wrap = True

pts5 = [
    ("📐 Reorder Point (ROP) Mathematical Formulation:", "Calculates exact reorder threshold: ROP = (d × L) + SS (where d = average daily demand, L = lead time, SS = safety stock)."),
    ("🔮 Stockout Risk Scoring:", "Classifies inventory health into 🟢 Healthy, 🟡 Reorder Warning, and 🔴 Critical Stockout to automate replenishment."),
    ("⏳ 45-Day Perishable Expiry Watchlist:", "Flags items approaching expiration 45 days prior, enabling promotional clearance and zero stock wastage."),
    ("💰 Purchase Rate Stock Valuation:", "Executive inventory valuation strictly calculated using Wholesale Purchase Cost: Total Stock Valuation = ∑ (purchase_price × stock).")
]

for h, b in pts5:
    p = tf5.add_paragraph()
    p.text = h
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p.space_before = Pt(12)
    
    p_b = tf5.add_paragraph()
    p_b.text = b
    p_b.font.size = Pt(14)
    p_b.font.color.rgb = DARK_TEXT
    p_b.space_before = Pt(4)

# ── SLIDE 6: HINGLISH AI ASSISTANT ────────────────────
slide6 = prs.slides.add_slide(blank_layout)
add_header(slide6, "AI Assistant Module", "Hinglish Conversational NLU Query Engine")

tb6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0))
tf6 = tb6.text_frame
tf6.word_wrap = True

pts6 = [
    ("🗣️ Natural Language Processing (Hinglish/English):", "Parses natural retail queries like 'Sugar ka rate Malegaon mein' or 'Atta ka stock kitna hai' in real time."),
    ("🎯 Target Commodity Filtering:", "Filters exact product and market matches without displaying unrelated items."),
    ("🤖 AI Product Card Response:", "Renders results in clean, structured HTML cards showing Rate, Market, Stock Level, Supplier, and Verification Timestamp."),
    ("📊 Integrated Plotly Visualizations:", "Generates dynamic price comparison line charts and margin distribution bar graphs.")
]

for h, b in pts6:
    p = tf6.add_paragraph()
    p.text = h
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_before = Pt(12)
    
    p_b = tf6.add_paragraph()
    p_b.text = b
    p_b.font.size = Pt(14)
    p_b.font.color.rgb = DARK_TEXT
    p_b.space_before = Pt(4)

# ── SLIDE 7: POS BILLING & CRM ───────────────────────
slide7 = prs.slides.add_slide(blank_layout)
add_header(slide7, "Point of Sale & CRM", "Ultra-Fast POS Billing & Customer Loyalty System")

tb7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0))
tf7 = tb7.text_frame
tf7.word_wrap = True

pts7 = [
    ("🧾 Barcode POS Billing (< 1 Second):", "Supports USB/Bluetooth barcode scanning, automatic GST tax computation, and custom line-item discounts."),
    ("📲 1-Click WhatsApp Receipts:", "Sends digital PDF/Text invoice receipts directly to customer WhatsApp numbers, eliminating paper cost."),
    ("👥 Customer CRM Ledger:", "Tracks individual customer order histories, total lifetime spend, and outstanding balances."),
    ("🎁 Loyalty Points Program:", "Automates loyalty reward points accumulation (1 point per ₹100 spend) to boost repeat customer retention.")
]

for h, b in pts7:
    p = tf7.add_paragraph()
    p.text = h
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.space_before = Pt(12)
    
    p_b = tf7.add_paragraph()
    p_b.text = b
    p_b.font.size = Pt(14)
    p_b.font.color.rgb = DARK_TEXT
    p_b.space_before = Pt(4)

# ── SLIDE 8: ARCHITECTURE & SECURITY ─────────────────
slide8 = prs.slides.add_slide(blank_layout)
add_header(slide8, "Technical Architecture", "Enterprise Security, Compliance & Data Access Layer")

tb8 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0))
tf8 = tb8.text_frame
tf8.word_wrap = True

pts8 = [
    ("🔐 Salted Password Hashing:", "PBKDF2 HMAC SHA-256 password hashing with 100,000 iterations for bulletproof authentication."),
    ("🛡️ Role-Based Access Control (RBAC):", "Strict authorization permissions across Admin, Store Manager, and Cashier Staff accounts."),
    ("🏛️ Data Access Layer (DAL):", "Thread-safe SQLite DatabaseManager with indexed SQL execution (idx_products_name, idx_bills_date)."),
    ("🌐 Digital India & UX4G Compliance:", "256-bit SSL encryption, WCAG 1.4.4 accessibility compliance, JSON-LD Schema.org structured data, and OpenGraph metadata.")
]

for h, b in pts8:
    p = tf8.add_paragraph()
    p.text = h
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p.space_before = Pt(12)
    
    p_b = tf8.add_paragraph()
    p_b.text = b
    p_b.font.size = Pt(14)
    p_b.font.color.rgb = DARK_TEXT
    p_b.space_before = Pt(4)

# ── SLIDE 9: COMMERCIAL IMPACT METRICS ──────────────
slide9 = prs.slides.add_slide(blank_layout)
add_header(slide9, "Value Proposition", "Commercial Business Impact & Operational Gains")

metrics = [
    ("12 - 18%", "Margin Expansion", "By procuring inventory using real-time APMC Mandi wholesale benchmark rates."),
    ("< 1 Sec", "POS Invoice Speed", "Ultra-fast customer checkout with instant 1-click WhatsApp digital receipts."),
    ("0%", "Stockout Rates", "Mathematical Reorder Point (ROP) automation keeping critical SKUs in stock."),
    ("100%", "Government Compliance", "Priority 1 Agmarknet / eNAM data integration with 256-bit SSL security.")
]

for i, (val, title, desc) in enumerate(metrics):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 5.9)
    top = Inches(1.8 + row * 2.5)
    
    box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.6), Inches(2.2))
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = RGBColor(226, 232, 240)
    
    tb = slide9.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.2), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = val
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = NAVY
    p2.space_before = Pt(4)
    
    p3 = tf.add_paragraph()
    p3.text = desc
    p3.font.size = Pt(12)
    p3.font.color.rgb = DARK_TEXT
    p3.space_before = Pt(6)

# ── SLIDE 10: CONCLUSION & CREDENTIALS ──────────────
slide10 = prs.slides.add_slide(blank_layout)
bg10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
bg10.fill.solid()
bg10.fill.fore_color.rgb = NAVY
bg10.line.color.rgb = NAVY

t_box10 = slide10.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.5))
tf10 = t_box10.text_frame
tf10.word_wrap = True

p1 = tf10.paragraphs[0]
p1.text = "🚀 RetailMind AI — Ready for Commercial Deployment"
p1.font.size = Pt(36)
p1.font.bold = True
p1.font.color.rgb = WHITE

p2 = tf10.add_paragraph()
p2.text = "Transforming Indian Retail & Supermarket Procurement with Artificial Intelligence"
p2.font.size = Pt(18)
p2.font.color.rgb = LIGHT_BLUE
p2.space_before = Pt(14)

p3 = tf10.add_paragraph()
p3.text = "• Live Application URL: https://retailmind-ai-by-suraj.streamlit.app/\n• Codebase Repository: https://github.com/suraj27255093/retail-mind-ai.git\n• Official Contact & Onboarding: contact@retailmind.ai"
p3.font.size = Pt(15)
p3.font.color.rgb = RGBColor(203, 213, 225)
p3.space_before = Pt(24)

p4 = tf10.add_paragraph()
p4.text = "Architected & Developed by: Suraj V. Shewale (शेवाळे पाटील)\nNashik & Pune, Maharashtra, India"
p4.font.size = Pt(16)
p4.font.bold = True
p4.font.color.rgb = EMERALD
p4.space_before = Pt(30)

output_path = "RetailMind_AI_Presentation.pptx"
prs.save(output_path)
print(f"POWERPOINT DECK GENERATED SUCCESSFULLY: {output_path}")
